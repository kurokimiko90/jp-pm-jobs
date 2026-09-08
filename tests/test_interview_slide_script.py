"""tools.interview_slide_script の抽出・PII・口語・事実ゲート。"""
from __future__ import annotations

import json
from pathlib import Path

from tools import interview_slide_script as script


def material(
    number: int = 1,
    text: str = "プロダクトマネージャー経験は9年です。",
    *,
    sample: bool = False,
    warnings: list[str] | None = None,
) -> script.SlideMaterial:
    return script.SlideMaterial(
        number=number,
        title=f"slide {number}",
        text=text,
        warnings=warnings or [],
        has_sample_disclaimer=sample,
    )


def generated_row(
    number: int = 1,
    text: str = (
        "私はプロダクトマネージャーとして9年間、開発に携わってきました。"
        "要件定義からリリースまで担当しています。"
        "事業成果まで意識して進めることが強みです。"
    ),
    *,
    warnings: list[str] | None = None,
) -> dict:
    return {
        "slide": number,
        "title": f"slide {number}",
        "seconds": 30,
        "key_message": "経験の一貫性",
        "script": text,
        "warnings": warnings or [],
    }


class TestSlideSpec:
    def test_range_and_dedup(self):
        assert script.parse_slide_spec("1-3,3,5", 6) == [1, 2, 3, 5]

    def test_out_of_range_is_rejected(self):
        try:
            script.parse_slide_spec("1-7", 6)
        except ValueError as exc:
            assert "存在しない" in str(exc)
        else:
            raise AssertionError("out-of-range slide must fail")


class TestSlideExtraction:
    def test_extracts_selected_pages_and_sample_warning(self, tmp_path: Path):
        from pptx import Presentation

        source = tmp_path / "source.pptx"
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        slide1.shapes.title.text = "自己紹介"
        slide1.shapes.add_textbox(0, 0, 100, 100).text = "PM経験は9年です。"
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "成果"
        slide2.shapes.add_textbox(0, 0, 100, 100).text = (
            "初の黒字化\n※ 内容・数値はサンプルです。実績に差し替えてください。"
        )
        prs.save(source)

        materials, total = script.extract_slide_materials(source, "1-2")
        assert total == 2
        assert [m.number for m in materials] == [1, 2]
        assert "PM経験は9年" in materials[0].text
        assert materials[1].has_sample_disclaimer
        assert any("口頭で主張しない" in w for w in materials[1].warnings)


class TestPii:
    def test_name_and_reading_are_tokenized_before_external_prompt(self, monkeypatch):
        from tools import deid, pii_gate

        monkeypatch.setattr(
            deid,
            "load_profile",
            lambda: {"identity": {"name_ja": "山田 太郎", "name_romaji": "Taro Yamada"}},
        )
        monkeypatch.setattr(
            deid,
            "load_resume_contact",
            lambda: {
                "name_ja": "山田 太郎",
                "name_romaji": "Taro Yamada",
                "email": "taro@example.com",
            },
        )
        monkeypatch.setattr(
            pii_gate,
            "scrub_for_external",
            lambda text: (text.replace("taro@example.com", "***"), ["email"] if "@" in text else []),
        )
        source = [
            script.SlideMaterial(
                1,
                "山田　太郎（ヤマダ タロウ）",
                "山田 太郎（ヤマダ タロウ）\ntaro@example.com",
            )
        ]

        safe, local_name, count = script.deidentify_materials(source)
        prompt_text = safe[0].title + safe[0].text
        assert "山田" not in prompt_text
        assert "ヤマダ" not in prompt_text
        assert "taro@example.com" not in prompt_text
        assert script.CANDIDATE_NAME_TOKEN in prompt_text
        assert local_name == "山田 太郎"
        assert count >= 2


class TestPrompt:
    def test_contains_user_requirements_and_detected_optimizations(self):
        prompt = script.build_prompt(
            [material(text="AI に資料整理を任せ、人が判断する。")],
            "positioning:\n  title: PM\n",
        )
        for expected in (
            "自然な日本語",
            "一枚につき主張は一つ",
            "「貴社」ではなく「御社」",
            "サンプル",
            "AIは情報整理や初稿",
            "最終責任",
            script.CANDIDATE_NAME_TOKEN,
        ):
            assert expected in prompt


class TestValidation:
    def test_valid_output_passes(self):
        materials = [material()]
        data = {"slides": [generated_row()]}
        assert script.validate_generated(data, materials, "経験は9年") == []

    def test_fabricated_number_is_rejected(self):
        materials = [material()]
        data = {"slides": [generated_row(text="プロダクトマネージャー経験は15年です。")]}
        findings = "\n".join(script.validate_generated(data, materials, "経験は9年"))
        assert "15" in findings

    def test_sample_claim_is_rejected_even_if_slide_contains_it(self):
        materials = [
            material(
                text="初の黒字化。500万円の新規収益。※ 数値はサンプルです。",
                sample=True,
            )
        ]
        data = {"slides": [generated_row(text="担当事業を黒字化し、収益を改善しました。")]}
        findings = "\n".join(script.validate_generated(data, materials, ""))
        assert "サンプルページ" in findings

    def test_ai_workflow_requires_human_judgment(self):
        materials = [material(text="AI に差分と矛盾を整理させる。")]
        bad = {"slides": [generated_row(text="AIで差分と矛盾を整理しました。")]}
        assert "判断・責任" in "\n".join(
            script.validate_generated(bad, materials, "")
        )

        good = {
            "slides": [
                generated_row(
                    text=(
                        "AIで差分を整理しました。"
                        "優先順位と最終判断は私が担当しました。"
                    )
                )
            ]
        }
        assert script.validate_generated(good, materials, "") == []

    def test_written_and_meta_language_are_rejected(self):
        materials = [material()]
        data = {
            "slides": [
                generated_row(
                    text="資料上では、当該プロジェクトに寄与したと答えます。"
                )
            ]
        }
        findings = "\n".join(script.validate_generated(data, materials, ""))
        assert "メタ表現" in findings
        assert "硬い書面語" in findings


class TestOutput:
    def test_generation_retries_failed_gate_and_merges_detected_warning(self, monkeypatch):
        import llm

        materials = [
            material(
                warnings=["数字は面接前に確認する。"],
            )
        ]
        responses = iter(
            [
                json.dumps(
                    {"slides": [generated_row(text="経験は15年です。")]},
                    ensure_ascii=False,
                ),
                json.dumps(
                    {
                        "slides": [
                            generated_row(
                                text=(
                                    f"{script.CANDIDATE_NAME_TOKEN}と申します。"
                                    "プロダクトマネージャーとして9年間、"
                                    "開発に携わってきました。"
                                )
                            )
                        ]
                    },
                    ensure_ascii=False,
                ),
            ]
        )
        prompts: list[str] = []

        def fake_call(prompt: str, **_kwargs):
            prompts.append(prompt)
            return next(responses)

        monkeypatch.setattr(llm, "call", fake_call)
        result = script.generate_with_llm(
            "base prompt",
            materials,
            "経験は9年",
        )
        assert len(prompts) == 2
        assert "素材で確認できない数字" in prompts[1]
        assert result["slides"][0]["warnings"] == ["数字は面接前に確認する。"]

    def test_markdown_contains_rules_warnings_and_script(self):
        data = {
            "slides": [
                {
                    **generated_row(),
                    "warnings": ["サンプル値は話さない。"],
                }
            ]
        }
        output = script.render_markdown("自己紹介", "1-6", data, redaction_count=2)
        assert "この台本に組み込んだ生成条件" in output
        assert "自動検出した改善点" in output
        assert "サンプル値は話さない" in output
        assert "外部送信前の PII 置換件数：2件" in output

    def test_user_layer_is_not_overwritten_without_force(self, tmp_path: Path):
        existing = tmp_path / "script.md"
        existing.write_text("user content", encoding="utf-8")
        try:
            script._assert_writable_outputs([existing], force=False)
        except FileExistsError as exc:
            assert "--force" in str(exc)
        else:
            raise AssertionError("existing output must be protected")
