"""面接スライド生成 — 母版 slide15/16 の会社特化欄位を JD 対応内容へ差し替え。

母版 interview/slides_master.pptx（17 枚、個人敘事固定）のうち、
slide 15「WHY ME＋入社後の展開」と slide 16 結び 1 文だけを LLM で生成した
テキストに置換する。版式・数字・実績はすべて母版側に固定（LLM を通らない）。

PII: prompt には tools.deid.build_deid_profile() の白名單のみ送出。
氏名・連絡先は母版にのみ存在し、LLM へは一切渡らない。

用法:
    python3 -m tools.interview_slides --job-id 123            # → output/prep 側へ
    python3 -m tools.interview_slides --job-id 123 --no-llm   # prompt 落地のみ
"""
import argparse
import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
MASTER = ROOT / "interview" / "slides_master.pptx"

# token -> (JSON key, 字数上限)。上限は母版原文 +2〜3 割で設定（超過は溢れリスク）
FIELDS: dict[str, tuple[str, int]] = {
    "{{WHYME}}":         ("whyme", 95),
    "{{SHORT_TITLE_1}}": ("short_title_1", 13),
    "{{SHORT_TITLE_2}}": ("short_title_2", 11),
    "{{SHORT_BODY}}":    ("short_body", 58),
    "{{MID_TITLE_1}}":   ("mid_title_1", 13),
    "{{MID_TITLE_2}}":   ("mid_title_2", 11),
    "{{MID_BODY}}":      ("mid_body", 58),
    "{{LONG_TITLE_1}}":  ("long_title_1", 13),
    "{{LONG_TITLE_2}}":  ("long_title_2", 11),
    "{{LONG_BODY}}":     ("long_body", 58),
    "{{CLOSING}}":       ("closing", 55),
}
_KEYS = [k for k, _ in FIELDS.values()]

JD_MAX_CHARS = 4000


def _build_prompt(job: dict) -> str:
    from tools.deid import build_deid_profile

    deid = build_deid_profile(compact=True)
    jd = (job.get("raw_jd") or "")[:JD_MAX_CHARS]
    gap = ""
    if job.get("gap_analysis"):
        try:
            g = json.loads(job["gap_analysis"]) if isinstance(job["gap_analysis"], str) else job["gap_analysis"]
            gap = json.dumps(
                {k: g[k] for k in ("strengths", "gaps", "recommend_reason") if k in g},
                ensure_ascii=False)
        except (json.JSONDecodeError, TypeError):
            pass

    limits = "\n".join(
        f"- {key}: {limit} 字以内" for _, (key, limit) in FIELDS.items())

    return f"""\
あなたは日本企業の PM 面接を熟知したキャリアコーチです。
応募者の面接用スライド（17 枚の自己紹介 deck）のうち、会社特化の 2 枚だけを
この求人向けに書き換えます。以下 11 欄位を JSON で出力してください。

厳守事項：
- 応募者プロフィールに記載された事実・数字のみ使用。捏造は一切禁止
- 各 body には必ずプロフィール由来の具体的な証拠（数字・プロジェクト名・成果）を 1 つ以上入れる
- JD の業務内容・必須要件に直接対応させる（汎用的な美辞麗句は不可）
- 全文日本語。title は体言止め（句点なし）、body は「。」で終える
- 面接でスクリーンに映す前提。口頭で補足するので、文字は簡潔に

# 対象求人
- 会社: {job.get('company') or '不明'}
- 職位: {job['title']}

# JD 全文（抜粋）
{jd}

# gap 分析（参考 — 強み/ギャップの既存評価）
{gap or '（なし）'}

# 応募者プロフィール（去識別化済み）
```yaml
{deid}
```

# 出力欄位（すべて必須）

1. **whyme** — 「なぜ私か」の総括 1〜2 文。原文の型「A × B × C を横断。〜できる PdM として貢献します。」を踏襲し、
   この会社のドメイン・フェーズに合わせて要素を選び直す
2. **short_title_1 / short_title_2** — 入社後すぐ（〜半年）に貢献できることのタイトル 2 行
   （例の型: 「決済・ポイント案件を」「即リード」）。short_body に根拠（実経験の数字）
3. **mid_title_1 / mid_title_2 / mid_body** — 1〜2 年での貢献。AI 活用の実運用化など、
   プロフィールの実績を土台にこの会社で再現すること
4. **long_title_1 / long_title_2 / long_body** — 中長期の役割像。この会社の事業・組織文脈に沿って
5. **closing** — deck 結びの 1 文。原文の型「この進め方は、御社の◯◯にも再現可能な形で展開できます。」の
   ◯◯をこの会社のプロダクト・事業に置換

# 字数上限（超過禁止 — スライド枠から溢れる）
{limits}

# 出力形式
JSON のみ。コードブロック・解説文は不要。
{{"whyme": "...", "short_title_1": "...", "short_title_2": "...", "short_body": "...",
 "mid_title_1": "...", "mid_title_2": "...", "mid_body": "...",
 "long_title_1": "...", "long_title_2": "...", "long_body": "...", "closing": "..."}}
"""


def _extract_json(raw: str) -> dict | None:
    m = re.search(r"\{.*\}", raw, re.DOTALL)
    if not m:
        return None
    try:
        data = json.loads(m.group(0))
    except json.JSONDecodeError:
        return None
    return data if isinstance(data, dict) else None


def _validate(data: dict) -> list[str]:
    """回傳問題清單（空 = 通過）。長度超上限 1.3 倍才視為硬錯誤。"""
    problems = []
    for _, (key, limit) in FIELDS.items():
        val = data.get(key)
        if not val or not str(val).strip():
            problems.append(f"{key}: 缺失或為空")
        elif "{{" in str(val):
            problems.append(f"{key}: 含未解決佔位符")
        elif len(str(val)) > limit * 1.3:
            problems.append(f"{key}: {len(str(val))} 字超過上限 {limit} 字")
    return problems


def build_fields(job: dict, no_llm: bool = False,
                 prompts_dir: Path | None = None) -> dict | None:
    """LLM 生成 11 欄位。失敗（或 no_llm）時 prompt 落地並回傳 None。"""
    prompt = _build_prompt(job)

    def _dump(reason: str) -> None:
        if prompts_dir:
            prompts_dir.mkdir(parents=True, exist_ok=True)
            (prompts_dir / "03_slides.prompt.md").write_text(prompt, encoding="utf-8")
            print(f"  {reason} → prompt 落地: {prompts_dir.name}/03_slides.prompt.md")

    if no_llm:
        _dump("→ --no-llm")
        return None

    from tools import miko_llm
    if not miko_llm.is_available():
        _dump("✗ 指揮中心不可用")
        return None

    accept = {"includesAll": [f'"{k}"' for k in _KEYS]}
    feedback = ""
    for attempt in (1, 2):
        try:
            raw = miko_llm.text(prompt + feedback, timeout=300, opts={"accept": accept})
        except Exception as e:
            _dump(f"✗ 指揮中心失敗（{e}）")
            return None
        data = _extract_json(raw or "")
        problems = _validate(data) if data else ["JSON 抽取失敗"]
        if not problems:
            return data
        if attempt == 1:
            feedback = ("\n\n# 前回出力の問題（修正して再出力）\n"
                        + "\n".join(f"- {p}" for p in problems))
            print(f"  ⚠ 驗收未過（{'; '.join(problems)}）→ 重試 1 次")
    _dump(f"✗ 重試後仍未過驗收（{'; '.join(problems)}）")
    return None


def render(fields: dict, out_path: Path) -> None:
    """母版佔位符替換 → out_path。任何 token 殘留即失敗。"""
    from pptx import Presentation

    if not MASTER.exists():
        raise FileNotFoundError(
            f"{MASTER} 不存在 — 先跑 python3 -m tools.oneoff.make_slides_master <deck.pptx>")
    prs = Presentation(str(MASTER))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "{{" not in run.text:
                        continue
                    for token, (key, _) in FIELDS.items():
                        if token in run.text:
                            run.text = run.text.replace(token, str(fields[key]).strip())
    leftovers = [
        run.text
        for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs
        if "{{" in run.text
    ]
    if leftovers:
        raise ValueError(f"佔位符殘留: {leftovers}")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def generate(job: dict, out_path: Path, no_llm: bool = False,
             prompts_dir: Path | None = None) -> bool:
    fields = build_fields(job, no_llm=no_llm, prompts_dir=prompts_dir)
    if fields is None:
        return False
    render(fields, out_path)
    # 欄位 JSON 併存 — 面接前微調時不用解 pptx
    out_path.with_suffix(".fields.json").write_text(
        json.dumps(fields, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"  ✓ {out_path.name}（slide15/16 特化 11 欄位）")
    return True


def main() -> None:
    p = argparse.ArgumentParser(description="面接スライド生成（母版差し替え）")
    p.add_argument("--job-id", type=int, required=True)
    p.add_argument("--out", default=None, help="出力 pptx パス（省略時 output/prep 配下）")
    p.add_argument("--no-llm", action="store_true")
    args = p.parse_args()

    from tracker.db import connect
    with connect() as conn:
        row = conn.execute("SELECT * FROM jobs WHERE id = ?", (args.job_id,)).fetchone()
    if row is None:
        sys.exit(f"job_id {args.job_id} が DB に見つかりません")
    job = dict(row)

    out = Path(args.out) if args.out else ROOT / "output" / "prep" / f"{args.job_id}_slides.pptx"
    ok = generate(job, out, no_llm=args.no_llm, prompts_dir=out.parent / "_prompts")
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
