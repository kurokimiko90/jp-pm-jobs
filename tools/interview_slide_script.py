"""PPTX から日本語の面接用スライド解説詞を生成する。

指定したスライドの文字を抽出し、候補者プロフィールの去識別化済み白名單と照合しながら、
日本の中途面接でそのまま話せる短い台本を生成する。元 PPTX は変更しない。

生成時に組み込む主な品質ルール:
- 一枚につき一つの主張、25〜40 秒、3〜5 文
- です・ます調の自然な話し言葉。投影片の読み上げや過剰敬語を避ける
- 「貴社」ではなく「御社」、「〜させていただきました」の多用を避ける
- 候補者資料にない経験・数字を作らない
- 「サンプル」「差し替え」等があるページの数値・成果を口頭主張しない
- AI は情報整理、人が優先順位・判断・最終責任を担うと明確にする

PII:
外部 LLM へ送る候補者プロフィールは tools.deid.build_deid_profile() の白名單のみ。
PPTX 由来テキストも tools.pii_gate.scrub_for_external() を通し、氏名はローカルの
<CANDIDATE_NAME> トークンへ置換してから送る。実名の復元は LLM 呼び出し後だけ行う。

用法:
    python3 -m tools.interview_slide_script \
      ~/Downloads/自己紹介.pptx --slides 1-6

    # LLM を呼ばず、去識別化済み prompt だけ確認
    python3 -m tools.interview_slide_script \
      ~/Downloads/自己紹介.pptx --slides 1-6 --no-llm
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

PROMPT_VERSION = 1
CANDIDATE_NAME_TOKEN = "<CANDIDATE_NAME>"
DEFAULT_SLIDES = "1-6"
MIN_SECONDS = 20
MAX_SECONDS = 45
MAX_SENTENCE_CHARS = 90
MAX_SCRIPT_CHARS = 320
MAX_RETRY = 2

SAMPLE_MARKERS = (
    "サンプル",
    "ダミー",
    "仮データ",
    "仮の数値",
    "差し替えて",
    "差し替え",
    "TODO",
    "要確認",
)
SAMPLE_CLAIM_RE = re.compile(
    r"(?:黒字化|赤字|利益|売上|収益|コスト|工数).{0,14}"
    r"(?:改善|向上|増加|削減|短縮|転換|達成)"
)
NUMBER_RE = re.compile(r"(?<![A-Za-z])\d+(?:[,.]\d+)*(?:\.\d+)?%?")
PLAIN_STYLE_RE = re.compile(r"(?:である|だ。|だった。)")
MARKDOWN_RE = re.compile(r"(?:\*\*|^#{1,6}\s|^\s*[-*]\s|```)", re.MULTILINE)

META_TERMS = (
    "資料上では",
    "面接では",
    "スライドには",
    "このスライド",
    "ここに書いて",
    "と答えます",
    "と話します",
    "と表現します",
)
STIFF_TERMS = (
    "当該",
    "上述",
    "下記",
    "勘案",
    "ご説明申し上げ",
    "推進してまいりました",
    "寄与しました",
    "責任分界",
    "精緻化条件",
)
HUMAN_JUDGMENT_TERMS = (
    "私が判断",
    "人が判断",
    "最終判断",
    "優先順位",
    "責任",
    "最終確認",
    "意思決定",
)

# ユーザー指定と実スライド監査から得た改善点を一箇所で管理する。
# prompt と最終 Markdown の両方に同じ一覧を出し、要件の漂移を防ぐ。
GENERATION_RULES: tuple[str, ...] = (
    "日本の中途採用面接で、そのまま口に出せる自然な日本語にする。",
    "簡潔な話し言葉のです・ます調にし、暗記した作文や投影片の読み上げにしない。",
    "一枚につき主張は一つ。結論を先に述べ、3〜5文、25〜40秒を目安にする。",
    "日本人の面接で自然な敬意は保つが、過剰敬語や「〜させていただきました」の連発を避ける。",
    "口頭では「貴社」ではなく「御社」を使う。",
    "投影片内の文章・項目・数字を全部読まず、意味と事業価値を一つだけ補足する。",
    "各ページを独立した説明にせず、経歴・強み・事例が一本のキャリアストーリーにつながるようにする。",
    "候補者資料にない経験・成果・数字は一切作らない。確認できない内容は台本に入れない。",
    "「サンプル」「差し替え」「要確認」等があるページは、該当する数字や成果を口頭で主張しない。",
    "AI活用を話す場合は、AIは情報整理や初稿、人は優先順位・判断・最終責任を担うと明確にする。",
    "「資料上では」「面接では〜と答えます」等のコーチ視点・メタ表現を候補者の発言に入れない。",
    "一文を短くし、専門用語や抽象語を並べず、担当・判断・成果が伝わる動詞を使う。",
    "ページ間の接続は必要な箇所だけ自然に入れ、前ページの内容を繰り返さない。",
)


@dataclass
class SlideMaterial:
    number: int
    title: str
    text: str
    warnings: list[str] = field(default_factory=list)
    has_sample_disclaimer: bool = False


@dataclass
class GeneratedSlide:
    slide: int
    title: str
    seconds: int
    key_message: str
    script: str
    warnings: list[str] = field(default_factory=list)


def parse_slide_spec(spec: str, total: int) -> list[int]:
    """``1-6,8,10-12`` 形式を 1-origin のページ番号へ展開する。"""
    if total < 1:
        raise ValueError("PPTX にスライドがありません")
    pages: list[int] = []
    for raw in spec.split(","):
        part = raw.strip()
        if not part:
            continue
        if "-" in part:
            bits = [x.strip() for x in part.split("-", 1)]
            if not all(x.isdigit() for x in bits):
                raise ValueError(f"不正なスライド範囲: {part!r}")
            start, end = map(int, bits)
            if start > end:
                raise ValueError(f"開始ページが終了ページより後です: {part!r}")
            pages.extend(range(start, end + 1))
        elif part.isdigit():
            pages.append(int(part))
        else:
            raise ValueError(f"不正なスライド指定: {part!r}")
    pages = list(dict.fromkeys(pages))
    if not pages:
        raise ValueError("スライド指定が空です")
    bad = [p for p in pages if p < 1 or p > total]
    if bad:
        raise ValueError(f"存在しないスライド番号: {bad}（全 {total} ページ）")
    return pages


def _shape_texts(shape: Any) -> Iterable[str]:
    """グループ・表を含む shape からテキストを再帰的に取り出す。"""
    if getattr(shape, "has_text_frame", False):
        text = (shape.text or "").strip()
        if text:
            yield text
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = (cell.text or "").strip()
                if text:
                    yield text
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is not None:
        for child in child_shapes:
            yield from _shape_texts(child)


def _dedupe_texts(texts: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for raw in texts:
        text = re.sub(r"[ \t]+", " ", raw).strip()
        if text and text not in seen:
            seen.add(text)
            result.append(text)
    return result


def detect_slide_warnings(text: str) -> tuple[list[str], bool]:
    """サンプル値・未解決 placeholder・過密ページを決定論的に検知する。"""
    warnings: list[str] = []
    matched = [marker for marker in SAMPLE_MARKERS if marker.lower() in text.lower()]
    has_sample = any(
        marker.lower() in text.lower()
        for marker in ("サンプル", "ダミー", "仮データ", "仮の数値", "差し替えて")
    )
    if has_sample:
        warnings.append(
            "サンプルまたは差し替え前の内容を検出。実績に置換するまで数値・成果を口頭で主張しない。"
        )
    unresolved = [m for m in matched if m.lower() in {"todo", "要確認", "差し替え"}]
    if unresolved:
        warnings.append(
            "未解決の注記を検出（" + "／".join(dict.fromkeys(unresolved)) + "）。面接前に確認する。"
        )
    if len(text) > 1200:
        warnings.append("情報量が多いページ。台本では一つの主張だけに絞る。")
    if len(NUMBER_RE.findall(text)) >= 12:
        warnings.append("数字が多いページ。数字を列挙せず、意味のある一つだけを使う。")
    return warnings, has_sample


def extract_slide_materials(pptx_path: Path, slide_spec: str) -> tuple[list[SlideMaterial], int]:
    """PPTX の指定ページを読み取り、視覚順に近い形でテキストをまとめる。"""
    from pptx import Presentation

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX が見つかりません: {pptx_path}")
    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError(f".pptx ファイルを指定してください: {pptx_path}")

    prs = Presentation(str(pptx_path))
    selected = parse_slide_spec(slide_spec, len(prs.slides))
    materials: list[SlideMaterial] = []
    for page in selected:
        slide = prs.slides[page - 1]
        ordered_shapes = sorted(
            slide.shapes,
            key=lambda sh: (int(getattr(sh, "top", 0)), int(getattr(sh, "left", 0))),
        )
        parts = _dedupe_texts(
            text for shape in ordered_shapes for text in _shape_texts(shape)
        )
        text = "\n".join(parts)
        title = ""
        try:
            title_shape = slide.shapes.title
            if title_shape is not None:
                title = (title_shape.text or "").strip()
        except (AttributeError, KeyError):
            pass
        if not title and parts:
            title = parts[0].splitlines()[0][:80]
        warnings, has_sample = detect_slide_warnings(text)
        if not text:
            warnings.append("抽出できる文字がないため、画像・図表の目視確認が必要。")
        materials.append(
            SlideMaterial(
                number=page,
                title=title or f"{page}ページ目",
                text=text,
                warnings=warnings,
                has_sample_disclaimer=has_sample,
            )
        )
    return materials, len(prs.slides)


def _flexible_name_pattern(name: str) -> re.Pattern[str]:
    escaped = re.escape(name.strip())
    escaped = escaped.replace(r"\ ", r"[\s\u3000]*")
    escaped = escaped.replace(r"\　", r"[\s\u3000]*")
    # 日本語名の後ろにある読み仮名も同時に外へ出さない。
    return re.compile(escaped + r"(?:\s*[（(][^）)\n]{1,40}[）)])?")


def deidentify_materials(
    materials: list[SlideMaterial],
) -> tuple[list[SlideMaterial], str, int]:
    """PPTX テキストを LLM 送信用に清洗し、ローカル復元用の氏名を返す。"""
    from tools.deid import load_profile, load_resume_contact
    from tools.pii_gate import scrub_for_external

    profile = load_profile()
    resume_contact = load_resume_contact()
    identity = profile.get("identity", {}) or {}
    names = [
        resume_contact.get("name_ja"),
        resume_contact.get("name_romaji"),
        identity.get("name_ja"),
        identity.get("name_romaji"),
    ]
    names = [str(v).strip() for v in names if isinstance(v, str) and v.strip()]
    local_name = str(
        resume_contact.get("name_ja") or identity.get("name_ja") or "本人"
    ).strip()

    redaction_count = 0
    sanitized: list[SlideMaterial] = []
    for material in materials:
        title = material.title
        text = material.text
        for name in dict.fromkeys(names):
            pattern = _flexible_name_pattern(name)
            title, count_title = pattern.subn(CANDIDATE_NAME_TOKEN, title)
            text, count_text = pattern.subn(CANDIDATE_NAME_TOKEN, text)
            redaction_count += count_title + count_text
        title, title_findings = scrub_for_external(title)
        text, text_findings = scrub_for_external(text)
        redaction_count += len(title_findings) + len(text_findings)
        sanitized.append(
            SlideMaterial(
                number=material.number,
                title=title,
                text=text,
                warnings=list(material.warnings),
                has_sample_disclaimer=material.has_sample_disclaimer,
            )
        )
    return sanitized, local_name, redaction_count


def _slide_blocks(materials: list[SlideMaterial]) -> str:
    blocks: list[str] = []
    for slide in materials:
        warning_text = "\n".join(f"- {w}" for w in slide.warnings) or "- なし"
        blocks.append(
            f"""\
## Slide {slide.number}
タイトル: {slide.title}
自動検出した注意:
{warning_text}
抽出テキスト:
<slide_text>
{slide.text}
</slide_text>"""
        )
    return "\n\n".join(blocks)


def build_prompt(materials: list[SlideMaterial], deid_profile: str) -> str:
    """ユーザー要件と改善ルールを全て含む、外部送信用 prompt を作る。"""
    rules = "\n".join(f"{idx}. {rule}" for idx, rule in enumerate(GENERATION_RULES, 1))
    slide_numbers = [m.number for m in materials]
    expected = ", ".join(map(str, slide_numbers))
    return f"""\
あなたは、日本企業の中途採用面接とプレゼンテーションに精通した日本語スピーチコーチです。
候補者が自己紹介スライドを使って面接で話すための、短く自然な解説詞を作成してください。

# コミュニケーションの目的
面接官が、候補者のキャリアの一貫性、PMとしての判断軸、再現可能な仕事の進め方を、
短時間で具体的に理解できるようにすること。

# 生成ルール（全て必須）
{rules}

# 事実と安全性
- 候補者プロフィールとスライド抽出テキストにある事実だけを使用する。
- 自動検出した注意は候補者が口に出す文章へ混ぜず、warnings に入れる。
- サンプル注記のあるページでは、例示の数字や成果を script に入れない。
- 候補者名が必要な箇所は、必ず文字列 {CANDIDATE_NAME_TOKEN} をそのまま使う。
- 去識別化された「本人」を候補者の一人称に使わず、「私」に言い換える。
- 外部知識、会社情報、推測による補完は禁止。

# 候補者プロフィール（去識別化済み白名單）
<candidate_profile>
{deid_profile}
</candidate_profile>

# 対象スライド
{_slide_blocks(materials)}

# 出力形式
JSON のみ。コードブロック、前置き、後書きは不要。
slides は必ず {expected} の順で全ページを一回ずつ含める。
seconds は {MIN_SECONDS}〜{MAX_SECONDS} の整数。
script は話す本文のみ。key_message はそのページで面接官に残す一つの要点。
warnings は「話さない注意」であり、script には含めない。

{{
  "slides": [
    {{
      "slide": {slide_numbers[0]},
      "title": "ページ内容に沿った短いタイトル",
      "seconds": 35,
      "key_message": "面接官に残す一つの要点",
      "script": "面接でそのまま話す自然な日本語。",
      "warnings": []
    }}
  ]
}}
"""


def _extract_json(raw: str) -> dict[str, Any] | None:
    start = raw.find("{")
    end = raw.rfind("}")
    if start < 0 or end <= start:
        return None
    try:
        data = json.loads(raw[start : end + 1])
    except json.JSONDecodeError:
        return None
    return data if isinstance(data, dict) else None


def _norm_number(token: str) -> str:
    return token.replace(",", "").rstrip("%")


def _known_numbers(text: str) -> set[str]:
    return {_norm_number(m.group()) for m in NUMBER_RE.finditer(text)}


def _unsupported_numbers(script: str, sources: str) -> list[str]:
    """列挙用 0〜10 を除き、素材に存在しない数値を返す。"""
    known = _known_numbers(sources)
    unsupported: set[str] = set()
    for match in NUMBER_RE.finditer(script):
        token = match.group()
        normalized = _norm_number(token)
        suffix = script[match.end() : match.end() + 5]
        if normalized in {str(n) for n in range(11)} and re.match(
            r"(?:つ|点|項目|ページ|つ目|点目|種類)", suffix
        ):
            continue
        if normalized not in known:
            unsupported.add(token)
    return sorted(unsupported)


def _sources_for_slide(
    slide: SlideMaterial,
    all_materials: list[SlideMaterial],
    deid_profile: str,
) -> str:
    """サンプルページ自身の数字を権威ソースから外す。"""
    trusted_slide_text = "\n".join(
        item.text for item in all_materials if not item.has_sample_disclaimer
    )
    if not slide.has_sample_disclaimer:
        trusted_slide_text += "\n" + slide.text
    return deid_profile + "\n" + trusted_slide_text


def _sentence_lengths(script: str) -> list[int]:
    return [
        len(sentence.strip())
        for sentence in re.split(r"[。！？\n]", script)
        if sentence.strip()
    ]


def _requires_ai_role_clarity(text: str) -> bool:
    """AI が作業主体として説明されるページだけ、判断責任の明示を必須にする。"""
    return bool(
        re.search(r"(?:AI|LLM)\s*(?:に|が|へ|を)\s*", text)
        or "AIにすべてを任せ" in text
        or "AIの役割" in text
    )


def validate_generated(
    data: dict[str, Any],
    materials: list[SlideMaterial],
    deid_profile: str,
) -> list[str]:
    """生成 JSON、口語、数字、サンプル主張、AI 責任分工を機械検収する。"""
    problems: list[str] = []
    rows = data.get("slides")
    if not isinstance(rows, list):
        return ["slides が配列ではありません"]

    expected = [m.number for m in materials]
    actual = [row.get("slide") for row in rows if isinstance(row, dict)]
    if actual != expected:
        problems.append(f"スライド順・網羅が不正: expected={expected}, actual={actual}")
        return problems

    by_number = {m.number: m for m in materials}
    for row in rows:
        page = row["slide"]
        prefix = f"slide {page}"
        for key in ("title", "key_message", "script"):
            if not isinstance(row.get(key), str) or not row[key].strip():
                problems.append(f"{prefix}: {key} が空")
        seconds = row.get("seconds")
        if not isinstance(seconds, int) or not MIN_SECONDS <= seconds <= MAX_SECONDS:
            problems.append(
                f"{prefix}: seconds は {MIN_SECONDS}〜{MAX_SECONDS} の整数にする"
            )
        if not isinstance(row.get("warnings", []), list):
            problems.append(f"{prefix}: warnings が配列ではない")

        script = str(row.get("script") or "")
        if "貴社" in script:
            problems.append(f"{prefix}: 面接口語で「貴社」を使用（「御社」にする）")
        if "本人" in script:
            problems.append(f"{prefix}: 去識別化語「本人」が残っている（「私」にする）")
        if PLAIN_STYLE_RE.search(script):
            problems.append(f"{prefix}: 常体・書面語が残っている")
        if MARKDOWN_RE.search(script):
            problems.append(f"{prefix}: script に Markdown 記号が残っている")
        meta = [term for term in META_TERMS if term in script]
        if meta:
            problems.append(f"{prefix}: メタ表現が残っている（{'／'.join(meta)}）")
        stiff = [term for term in STIFF_TERMS if term in script]
        if stiff:
            problems.append(f"{prefix}: 硬い書面語が残っている（{'／'.join(stiff)}）")
        if script.count("させていただ") > 1:
            problems.append(f"{prefix}: 「させていただく」が多すぎる")
        sentence_lengths = _sentence_lengths(script)
        if not 2 <= len(sentence_lengths) <= 6:
            problems.append(f"{prefix}: script は2〜6文にする")
        if len(script) > MAX_SCRIPT_CHARS:
            problems.append(f"{prefix}: script が{MAX_SCRIPT_CHARS}字を超えている")
        if max(sentence_lengths or [0]) > MAX_SENTENCE_CHARS:
            problems.append(f"{prefix}: {MAX_SENTENCE_CHARS}字を超える一文がある")

        material = by_number[page]
        unsupported = _unsupported_numbers(
            script, _sources_for_slide(material, materials, deid_profile)
        )
        if unsupported:
            problems.append(
                f"{prefix}: 素材で確認できない数字（{'／'.join(unsupported)}）"
            )
        if material.has_sample_disclaimer and SAMPLE_CLAIM_RE.search(script):
            problems.append(f"{prefix}: サンプルページの成果を実績として話している")
        if _requires_ai_role_clarity(material.text) and "AI" in script and not any(
            term in script for term in HUMAN_JUDGMENT_TERMS
        ):
            problems.append(f"{prefix}: AIと人の判断・責任の分担が不明確")
    return problems


def _merge_warnings(data: dict[str, Any], materials: list[SlideMaterial]) -> None:
    detected = {m.number: m.warnings for m in materials}
    for row in data.get("slides", []):
        model_warnings = row.get("warnings") if isinstance(row.get("warnings"), list) else []
        row["warnings"] = list(
            dict.fromkeys([*detected.get(row.get("slide"), []), *map(str, model_warnings)])
        )


def _restore_tokens(value: Any, local_name: str) -> Any:
    if isinstance(value, str):
        return value.replace(CANDIDATE_NAME_TOKEN, local_name)
    if isinstance(value, list):
        return [_restore_tokens(item, local_name) for item in value]
    if isinstance(value, dict):
        return {key: _restore_tokens(item, local_name) for key, item in value.items()}
    return value


def generate_with_llm(
    prompt: str,
    materials: list[SlideMaterial],
    deid_profile: str,
    timeout: int = 420,
) -> dict[str, Any]:
    """最大 2 回。違反があれば機械 feedback を付けて一度だけ修復する。"""
    from llm import call

    feedback = ""
    last_problems: list[str] = []
    for attempt in range(1, MAX_RETRY + 1):
        raw = call(
            prompt + feedback,
            timeout=timeout,
            accept={"includesAll": ['"slides"', '"script"', '"key_message"']},
        )
        data = _extract_json(raw)
        last_problems = (
            validate_generated(data, materials, deid_profile)
            if data is not None
            else ["JSON を抽出できません"]
        )
        if data is not None and not last_problems:
            _merge_warnings(data, materials)
            data["prompt_version"] = PROMPT_VERSION
            data["generation_rules"] = list(GENERATION_RULES)
            return data
        if attempt < MAX_RETRY:
            feedback = (
                "\n\n# 前回出力の問題（全て修正し、JSON 全体を再出力）\n"
                + "\n".join(f"- {problem}" for problem in last_problems[:30])
            )
    raise ValueError("生成結果が品質検査を通過しません: " + "; ".join(last_problems))


def render_markdown(
    source_name: str,
    selected_spec: str,
    data: dict[str, Any],
    redaction_count: int,
) -> str:
    """台本・生成条件・自動検出した改善点を一つの Markdown にまとめる。"""
    rows = [
        f"# {source_name}｜面接用スライド解説詞",
        "",
        f"対象：{selected_spec}ページ｜Prompt version：{PROMPT_VERSION}",
        "",
        "## この台本に組み込んだ生成条件",
        "",
    ]
    rows += [f"- {rule}" for rule in GENERATION_RULES]
    rows += [
        "",
        "## スライド別台本",
        "",
    ]
    for row in data["slides"]:
        rows += [
            f"### {row['slide']}ページ目｜{row['title']}（約{row['seconds']}秒）",
            "",
            f"**このページで残す要点：** {row['key_message']}",
            "",
            row["script"].strip(),
            "",
        ]
        for warning in row.get("warnings", []):
            rows.append(f"> 面接前の注意：{warning}")
        if row.get("warnings"):
            rows.append("")

    all_warnings = [
        (row["slide"], warning)
        for row in data["slides"]
        for warning in row.get("warnings", [])
    ]
    rows += [
        "## 自動検出した改善点",
        "",
    ]
    if all_warnings:
        rows += [f"- {page}ページ目：{warning}" for page, warning in all_warnings]
    else:
        rows.append("- サンプル値・未解決注記などの明確なリスクは検出されませんでした。")
    rows += [
        "",
        "## 面接直前の使い方",
        "",
        "- 台本を丸暗記せず、各ページの「残す要点」だけを覚える。",
        "- スライドではなく面接官を見て話し、数字を言う前に根拠を再確認する。",
        "- AIの事例では、最後に自分の判断・確認・責任を一言添える。",
        "- ページを切り替えたら一拍置き、面接官の反応を確認する。",
        "",
        f"外部送信前の PII 置換件数：{redaction_count}件（実名の復元はローカル出力時のみ）",
        "",
    ]
    return "\n".join(rows)


def _assert_writable_outputs(paths: list[Path], force: bool) -> None:
    existing = [path for path in paths if path.exists()]
    if existing and not force:
        joined = ", ".join(str(path) for path in existing)
        raise FileExistsError(
            f"既存の User Layer 出力を上書きしません: {joined}（必要なら --force）"
        )


def run(
    pptx_path: Path,
    slide_spec: str,
    out_path: Path,
    *,
    no_llm: bool = False,
    force: bool = False,
    timeout: int = 420,
) -> tuple[Path, Path | None]:
    """CLI 本体。成功時 (Markdown または prompt, JSON または None) を返す。"""
    from tools.deid import build_deid_profile

    materials, _ = extract_slide_materials(pptx_path, slide_spec)
    safe_materials, local_name, redaction_count = deidentify_materials(materials)
    deid_profile = build_deid_profile(compact=True)
    prompt = build_prompt(safe_materials, deid_profile)

    if no_llm:
        prompt_path = out_path.with_suffix(".prompt.txt")
        _assert_writable_outputs([prompt_path], force)
        prompt_path.parent.mkdir(parents=True, exist_ok=True)
        prompt_path.write_text(prompt, encoding="utf-8")
        return prompt_path, None

    json_path = out_path.with_suffix(".json")
    _assert_writable_outputs([out_path, json_path], force)
    try:
        data = generate_with_llm(prompt, safe_materials, deid_profile, timeout=timeout)
    except (RuntimeError, ValueError):
        prompt_path = out_path.with_suffix(".prompt.txt")
        if force or not prompt_path.exists():
            prompt_path.parent.mkdir(parents=True, exist_ok=True)
            prompt_path.write_text(prompt, encoding="utf-8")
        raise

    restored = _restore_tokens(data, local_name)
    output_payload = {
        "source": pptx_path.name,
        "selected_slides": [m.number for m in materials],
        "redaction_count": redaction_count,
        **restored,
    }
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        render_markdown(
            pptx_path.stem,
            slide_spec,
            output_payload,
            redaction_count,
        ),
        encoding="utf-8",
    )
    json_path.write_text(
        json.dumps(output_payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return out_path, json_path


def main() -> None:
    parser = argparse.ArgumentParser(
        description="PPTX から日本語の面接用スライド解説詞を生成"
    )
    parser.add_argument("pptx", help="入力 .pptx")
    parser.add_argument(
        "--slides",
        default=DEFAULT_SLIDES,
        help="対象ページ。例: 1-6,8,10-12（default: 1-6）",
    )
    parser.add_argument(
        "--out",
        default=None,
        help="出力 Markdown。省略時は output/<PPT名>_面接用解説詞.md",
    )
    parser.add_argument(
        "--no-llm",
        action="store_true",
        help="LLM を呼ばず、去識別化済み prompt.txt だけ出力",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="既存の出力を明示的に上書き",
    )
    parser.add_argument("--timeout", type=int, default=420)
    args = parser.parse_args()

    pptx_path = Path(args.pptx).expanduser().resolve()
    out_path = (
        Path(args.out).expanduser().resolve()
        if args.out
        else ROOT / "output" / f"{pptx_path.stem}_面接用解説詞.md"
    )
    try:
        primary, sidecar = run(
            pptx_path,
            args.slides,
            out_path,
            no_llm=args.no_llm,
            force=args.force,
            timeout=args.timeout,
        )
    except (FileNotFoundError, FileExistsError, RuntimeError, ValueError) as exc:
        sys.exit(f"✗ {exc}")

    if args.no_llm:
        print(f"✓ 去識別化済み prompt: {primary}")
    else:
        print(f"✓ 面接用解説詞: {primary}")
        print(f"✓ 構造化データ: {sidecar}")


if __name__ == "__main__":
    main()
